"""Document providers for fetching content from various URI schemes."""

import re
import tempfile
from pathlib import Path

import requests

from .base import Document, DocumentProvider, get_registry


# RFC 822 email header detection — matches common email headers at line start
_EMAIL_HEADER_RE = re.compile(
    r'^(From|To|Subject|Date|Message-ID|Cc|Bcc|Reply-To|Received'
    r'|MIME-Version|Content-Type|Return-Path|Delivered-To):\s',
    re.IGNORECASE | re.MULTILINE,
)

# Base64 block: 4+ consecutive lines of 60+ chars with only base64 alphabet
_BASE64_BLOCK_RE = re.compile(
    r'(?:^[A-Za-z0-9+/=]{60,}\n){4,}',
    re.MULTILINE,
)


def _strip_base64_blocks(text: str) -> str:
    """Remove base64-encoded blocks from text content.

    Emails with embedded MIME attachments (especially in corpora where
    multipart structure is flattened to text/plain) contain large blocks
    of base64. These are noise for summarization and embedding.
    """
    stripped = _BASE64_BLOCK_RE.sub('[base64 data removed]\n', text)
    return stripped


# Headers to skip — transport/encoding headers that aren't useful as tags
_EMAIL_SKIP_HEADERS = frozenset({
    'mime-version', 'content-type', 'content-transfer-encoding',
    'content-disposition', 'received', 'return-path', 'delivered-to',
    'x-mailer', 'x-originating-ip',
})


def extract_html_text(html_content: str) -> str:
    """Extract readable text from HTML, removing scripts and styles.

    Used by both FileDocumentProvider and HttpDocumentProvider to ensure
    consistent content regularization for embedding and summarization.

    Args:
        html_content: Raw HTML string

    Returns:
        Extracted text with whitespace normalized

    Raises:
        ImportError: If beautifulsoup4 is not installed
    """
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html_content, "html.parser")

    # Remove script and style elements
    for script in soup(["script", "style"]):
        script.decompose()

    # Get text
    text = soup.get_text()

    # Clean up whitespace
    lines = (line.strip() for line in text.splitlines())
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    return '\n'.join(chunk for chunk in chunks if chunk)


class FileDocumentProvider:
    """Fetches documents from the local filesystem.

    Supports file:// URIs and attempts to detect content type from extension.
    Performs text extraction for PDF and HTML files.
    """

    # Default max file size: 100MB
    MAX_FILE_SIZE = 100_000_000

    def __init__(self, max_size: int | None = None):
        self.max_size = max_size or self.MAX_FILE_SIZE

    EXTENSION_TYPES = {
        ".md": "text/markdown",
        ".markdown": "text/markdown",
        ".txt": "text/plain",
        ".py": "text/x-python",
        ".js": "text/javascript",
        ".ts": "text/typescript",
        ".json": "text/json",
        ".yaml": "text/yaml",
        ".yml": "text/yaml",
        ".html": "text/html",
        ".htm": "text/html",
        ".eml": "message/rfc822",
        ".mbox": "message/rfc822",
        ".css": "text/css",
        ".xml": "text/xml",
        ".rst": "text/x-rst",
        ".pdf": "application/pdf",
        # Office documents
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        # Apple iWork
        ".pages": "application/vnd.apple.pages",
        ".key": "application/vnd.apple.keynote",
        ".numbers": "application/vnd.apple.numbers",
        # Audio
        ".mp3": "audio/mpeg",
        ".flac": "audio/flac",
        ".aiff": "audio/aiff",
        ".aif": "audio/aiff",
        ".ogg": "audio/ogg",
        ".wav": "audio/wav",
        ".m4a": "audio/mp4",
        ".alac": "audio/mp4",
        ".wma": "audio/x-ms-wma",
        # Images
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".tiff": "image/tiff",
        ".tif": "image/tiff",
        ".webp": "image/webp",
        ".svg": "image/svg+xml",
        # Archives / binary (skip text extraction)
        ".zip": "application/zip",
        ".gz": "application/gzip",
        ".tar": "application/x-tar",
        ".7z": "application/x-7z-compressed",
        ".rar": "application/x-rar-compressed",
        ".bz2": "application/x-bzip2",
        ".dmg": "application/x-apple-diskimage",
        ".iso": "application/x-iso9660-image",
        ".xls": "application/vnd.ms-excel",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }
    
    def supports(self, uri: str) -> bool:
        """Check if this is a file:// URI or bare path."""
        return uri.startswith("file://") or uri.startswith("/")
    
    def fetch(self, uri: str) -> Document:
        """Read file content from the filesystem with text extraction for PDF/HTML."""
        # Normalize to path
        if uri.startswith("file://"):
            path_str = uri.removeprefix("file://")
        else:
            path_str = uri

        path = Path(path_str).resolve()

        if not path.exists():
            raise IOError(f"File not found: {path}")

        if not path.is_file():
            raise IOError(f"Not a file: {path}")

        # Reject paths outside user's home directory as a safety boundary
        from ..paths import validate_path_within_home

        try:
            validate_path_within_home(path)
        except ValueError:
            raise IOError(f"Path traversal blocked: {path} is outside home directory")

        # Check file size before processing
        file_size = path.stat().st_size
        if file_size > self.max_size:
            raise IOError(
                f"File too large: {file_size:,} bytes "
                f"(limit: {self.max_size:,} bytes). "
                f"Configure max_file_size in store config to increase."
            )

        # Detect content type
        suffix = path.suffix.lower()
        content_type = self.EXTENSION_TYPES.get(suffix, "text/plain")

        # Extract text based on file type.
        # Extraction failures are non-fatal: fall back to filename-only content
        # so that directory puts don't fail on one bad file.
        extracted_tags: dict[str, str] | None = None
        ocr_pages: list[int] | None = None
        email_attachments: list[dict] | None = None
        try:
            if suffix == ".pdf":
                content, ocr_pages = self._extract_pdf_text(path)
                if not content and ocr_pages:
                    # Fully scanned PDF — placeholder until background OCR
                    content = (
                        f"[Scanned document: {path.name}, "
                        f"{len(ocr_pages)} pages pending OCR]"
                    )
            elif suffix in (".html", ".htm"):
                content = self._extract_html_text(path)
            elif suffix in (".docx",):
                content, extracted_tags = self._extract_docx(path)
            elif suffix in (".pptx",):
                content, extracted_tags = self._extract_pptx(path)
            elif content_type and content_type.startswith("application/vnd.apple."):
                content, extracted_tags = self._extract_iwork_metadata(path, content_type)
            elif suffix == ".svg":
                content, extracted_tags = self._extract_svg_text(path)
            elif content_type and content_type.startswith("audio/"):
                content, extracted_tags = self._extract_audio_metadata(path)
            elif content_type and content_type.startswith("image/"):
                content, extracted_tags = self._extract_image_metadata(path)
                # Signal that this image should be OCR'd in background
                ocr_pages = [0]
            elif content_type == "message/rfc822":
                content, extracted_tags, email_attachments = self._extract_email(path)
            elif content_type and not content_type.startswith("text/"):
                # Binary file without a dedicated extractor — use filename only
                content = f"[{path.name}]"
            else:
                # Read as plain text
                content = path.read_text(encoding="utf-8")
                # Sniff for email format in extensionless files
                if self._detect_email(content):
                    content, extracted_tags, email_attachments = self._extract_email(path)
                    content_type = "message/rfc822"
        except Exception as e:
            import logging
            logging.getLogger(__name__).warning(
                "Text extraction failed for %s: %s", path.name, e
            )
            content = f"[{path.name}]"
            # Still try to get metadata tags from the file
            extracted_tags = self._extract_metadata_only(path, suffix)

        # Gather metadata
        stat = path.stat()
        metadata = {
            "size": stat.st_size,
            "modified": stat.st_mtime,
            "name": path.name,
        }
        # File creation time (birthtime on macOS, may exist on Linux 3.12+)
        try:
            metadata["birthtime"] = stat.st_birthtime
        except AttributeError:
            pass

        # Signal pages needing OCR to the caller (for background processing)
        if ocr_pages:
            metadata["_ocr_pages"] = ocr_pages

        # Signal email attachments to the caller (for child item creation)
        if email_attachments:
            metadata["_attachments"] = email_attachments

        return Document(
            uri=f"file://{path.resolve()}",  # Normalize to absolute
            content=content,
            content_type=content_type,
            metadata=metadata,
            tags=extracted_tags,
        )

    def _extract_metadata_only(self, path: Path, suffix: str) -> dict[str, str] | None:
        """Best-effort metadata extraction when content extraction fails.

        Tries to pull author/title from document properties without
        reading the full content. Returns None if nothing found.
        """
        tags: dict[str, str] = {}
        try:
            if suffix == ".pdf":
                from pypdf import PdfReader
                reader = PdfReader(path)
                info = reader.metadata
                if info:
                    if info.author:
                        tags["author"] = info.author
                    if info.title:
                        tags["title"] = info.title
            elif suffix == ".docx":
                from docx import Document as DocxDocument
                doc = DocxDocument(path)
                props = doc.core_properties
                if props.author:
                    tags["author"] = props.author
                if props.title:
                    tags["title"] = props.title
            elif suffix == ".pptx":
                from pptx import Presentation
                prs = Presentation(path)
                props = prs.core_properties
                if props.author:
                    tags["author"] = props.author
                if props.title:
                    tags["title"] = props.title
        except Exception:
            pass  # Metadata extraction is best-effort
        return tags or None

    def _extract_docx(self, path: Path) -> tuple[str, dict[str, str]]:
        """Extract text and metadata from DOCX file."""
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise IOError(
                f"DOCX support requires 'python-docx' library. "
                f"Install with: pip install python-docx\n"
                f"Cannot read DOCX: {path}"
            )

        try:
            doc = DocxDocument(path)
            parts = []

            # Extract paragraph text
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)

            # Extract table text
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))

            if not parts:
                raise IOError(f"No text extracted from DOCX: {path}")

            content = "\n\n".join(parts)

            # Extract metadata tags
            tags: dict[str, str] = {}
            props = doc.core_properties
            if props.author:
                tags["author"] = props.author
            if props.title:
                tags["title"] = props.title

            return content, tags or None
        except ImportError:
            raise
        except IOError:
            raise
        except Exception as e:
            raise IOError(f"Failed to extract text from DOCX {path}: {e}")

    def _extract_pptx(self, path: Path) -> tuple[str, dict[str, str]]:
        """Extract text and metadata from PPTX file."""
        try:
            from pptx import Presentation
        except ImportError:
            raise IOError(
                f"PPTX support requires 'python-pptx' library. "
                f"Install with: pip install python-pptx\n"
                f"Cannot read PPTX: {path}"
            )

        try:
            prs = Presentation(path)
            parts = []

            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)
                if slide_texts:
                    parts.append(f"Slide {i}:\n" + "\n".join(slide_texts))

                # Extract notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"Notes:\n{notes}")

            if not parts:
                raise IOError(f"No text extracted from PPTX: {path}")

            content = "\n\n".join(parts)

            # Extract metadata tags
            tags: dict[str, str] = {}
            props = prs.core_properties
            if props.author:
                tags["author"] = props.author
            if props.title:
                tags["title"] = props.title

            return content, tags or None
        except ImportError:
            raise
        except IOError:
            raise
        except Exception as e:
            raise IOError(f"Failed to extract text from PPTX {path}: {e}")

    def _extract_iwork_metadata(self, path: Path, content_type: str) -> tuple[str, dict[str, str]]:
        """Extract basic metadata from Apple iWork files (.pages, .key, .numbers).

        iWork files use a proprietary protobuf format (IWA) inside a ZIP archive.
        No production Python parser exists, so we extract what we can from the
        ZIP structure and return a non-text content type for media description.
        """
        import zipfile

        format_names = {
            "application/vnd.apple.pages": "Pages",
            "application/vnd.apple.keynote": "Keynote",
            "application/vnd.apple.numbers": "Numbers",
        }
        fmt = format_names.get(content_type, "iWork")
        tags: dict[str, str] = {"format": fmt.lower()}

        parts = [f"Apple {fmt} document: {path.name}"]

        try:
            if zipfile.is_zipfile(path):
                with zipfile.ZipFile(path, "r") as zf:
                    names = zf.namelist()
                    # Count embedded images
                    images = [n for n in names if n.startswith("Data/") and not n.endswith("/")]
                    if images:
                        parts.append(f"Contains {len(images)} embedded resource(s)")
        except Exception:
            pass  # ZIP inspection is best-effort

        return "\n".join(parts), tags

    def _extract_svg_text(self, path: Path) -> tuple[str, dict[str, str]]:
        """Extract meaningful text from SVG files.

        Pulls <title>, <desc>, and <text> elements from the XML,
        ignoring the visual markup (paths, shapes, transforms).
        """
        import xml.etree.ElementTree as ET

        try:
            tree = ET.parse(path)
        except ET.ParseError as e:
            raise IOError(f"Failed to parse SVG {path}: {e}")

        root = tree.getroot()
        # SVG namespace
        ns = {"svg": "http://www.w3.org/2000/svg"}

        parts = []
        tags: dict[str, str] = {}

        # Extract <title>
        for title in root.iter("{http://www.w3.org/2000/svg}title"):
            text = (title.text or "").strip()
            if text:
                parts.append(text)
                tags.setdefault("title", text)
        # Also check for non-namespaced title
        for title in root.iter("title"):
            text = (title.text or "").strip()
            if text and text not in parts:
                parts.append(text)
                tags.setdefault("title", text)

        # Extract <desc>
        for desc in root.iter("{http://www.w3.org/2000/svg}desc"):
            text = (desc.text or "").strip()
            if text:
                parts.append(text)
        for desc in root.iter("desc"):
            text = (desc.text or "").strip()
            if text and text not in parts:
                parts.append(text)

        # Extract <text> elements (visible text in the graphic)
        text_parts = []
        for elem in root.iter("{http://www.w3.org/2000/svg}text"):
            text = "".join(elem.itertext()).strip()
            if text:
                text_parts.append(text)
        for elem in root.iter("text"):
            text = "".join(elem.itertext()).strip()
            if text and text not in text_parts:
                text_parts.append(text)

        if text_parts:
            parts.append("\n".join(text_parts))

        if not parts:
            parts.append(f"SVG image: {path.name}")

        return "\n\n".join(parts), tags or None

    def _extract_audio_metadata(self, path: Path) -> tuple[str, dict[str, str]]:
        """Extract metadata from audio file as structured text."""
        try:
            from tinytag import TinyTag
        except ImportError:
            raise IOError(
                f"Audio metadata support requires 'tinytag' library. "
                f"Install with: pip install tinytag\n"
                f"Cannot read audio metadata: {path}"
            )

        try:
            tag = TinyTag.get(path)
            lines = []
            tags: dict[str, str] = {}

            if tag.title:
                lines.append(f"Title: {tag.title}")
                tags["title"] = tag.title
            if tag.artist:
                lines.append(f"Artist: {tag.artist}")
                tags["artist"] = tag.artist
            if tag.album:
                lines.append(f"Album: {tag.album}")
                tags["album"] = tag.album
            if tag.albumartist:
                lines.append(f"Album Artist: {tag.albumartist}")
            if tag.track:
                lines.append(f"Track: {tag.track}")
            if tag.genre:
                lines.append(f"Genre: {tag.genre}")
                tags["genre"] = tag.genre
            if tag.year:
                lines.append(f"Year: {tag.year}")
                tags["year"] = str(tag.year)
            if tag.duration:
                mins, secs = divmod(int(tag.duration), 60)
                lines.append(f"Duration: {mins}:{secs:02d}")
            if tag.bitrate:
                lines.append(f"Bitrate: {int(tag.bitrate)} kbps")
            if tag.samplerate:
                lines.append(f"Sample Rate: {tag.samplerate} Hz")
            if tag.comment:
                lines.append(f"Comment: {tag.comment}")

            if not lines:
                lines.append(f"Audio file: {path.name}")

            content = "\n".join(lines)
            return content, tags or None
        except ImportError:
            raise
        except Exception as e:
            raise IOError(f"Failed to extract audio metadata from {path}: {e}")

    def _extract_image_metadata(self, path: Path) -> tuple[str, dict[str, str]]:
        """Extract EXIF metadata from image file as structured text."""
        try:
            from PIL import Image
            from PIL.ExifTags import TAGS as EXIF_TAGS
        except ImportError:
            raise IOError(
                f"Image metadata support requires 'Pillow' library. "
                f"Install with: pip install Pillow\n"
                f"Cannot read image metadata: {path}"
            )

        # Guard against decompression bombs — allow large images (up to ~15800x15800)
        # but reject pathological ones that would exhaust memory.
        Image.MAX_IMAGE_PIXELS = 250_000_000

        try:
            img = Image.open(path)
            lines = []
            tags: dict[str, str] = {}

            # Basic image info
            w, h = img.size
            lines.append(f"Dimensions: {w}x{h}")
            tags["dimensions"] = f"{w}x{h}"
            lines.append(f"Format: {img.format}")

            # EXIF data
            exif = img.getexif()
            if exif:
                # Camera model
                model = exif.get(0x0110)  # Model
                make = exif.get(0x010F)  # Make
                if model:
                    camera = f"{make} {model}".strip() if make else model
                    lines.append(f"Camera: {camera}")
                    tags["camera"] = camera

                # Date taken
                date = exif.get(0x9003) or exif.get(0x0132)  # DateTimeOriginal or DateTime
                if date:
                    lines.append(f"Date: {date}")
                    tags["date"] = date

                # Focal length
                focal = exif.get(0x920A)  # FocalLength
                if focal:
                    if hasattr(focal, 'numerator'):
                        focal_val = focal.numerator / focal.denominator
                        lines.append(f"Focal Length: {focal_val:.0f}mm")
                    else:
                        lines.append(f"Focal Length: {focal}mm")

                # ISO
                iso = exif.get(0x8827)  # ISOSpeedRatings
                if iso:
                    lines.append(f"ISO: {iso}")

                # Exposure time
                exposure = exif.get(0x829A)  # ExposureTime
                if exposure:
                    if hasattr(exposure, 'numerator'):
                        if exposure.numerator == 1:
                            lines.append(f"Exposure: 1/{exposure.denominator}s")
                        else:
                            lines.append(f"Exposure: {exposure.numerator}/{exposure.denominator}s")
                    else:
                        lines.append(f"Exposure: {exposure}s")

                # F-number
                fnumber = exif.get(0x829D)  # FNumber
                if fnumber:
                    if hasattr(fnumber, 'numerator'):
                        fval = fnumber.numerator / fnumber.denominator
                        lines.append(f"Aperture: f/{fval:.1f}")
                    else:
                        lines.append(f"Aperture: f/{fnumber}")

            img.close()

            if not lines:
                lines.append(f"Image file: {path.name}")

            content = "\n".join(lines)
            return content, tags or None
        except ImportError:
            raise
        except Exception as e:
            raise IOError(f"Failed to extract image metadata from {path}: {e}")

    def _extract_pdf_text(self, path: Path) -> tuple[str, list[int]]:
        """Extract text from PDF file, reporting pages that need OCR.

        Returns:
            Tuple of (extracted_text, ocr_needed_page_indices).
            Text is from pages with embedded text layers only.
            OCR pages are deferred to background processing.
        """
        try:
            from pypdf import PdfReader
        except ImportError:
            raise IOError(
                f"PDF support requires 'pypdf' library. "
                f"Install with: pip install pypdf\n"
                f"Cannot read PDF: {path}"
            )

        try:
            reader = PdfReader(path)
            text_parts: list[tuple[int, str]] = []
            ocr_needed: list[int] = []

            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text and text.strip():
                    text_parts.append((i, text))
                else:
                    ocr_needed.append(i)

            if text_parts:
                text_parts.sort(key=lambda t: t[0])
                content = "\n\n".join(text for _, text in text_parts)
                # Check if extracted text is low-quality (schematics, scanned
                # docs with minimal embedded text, garbled extraction).
                # If so, flag all pages for OCR as well — text is kept as
                # fallback but OCR results will replace it.
                if self._is_low_quality_text(content, len(reader.pages)):
                    ocr_needed = list(range(len(reader.pages)))
            elif ocr_needed:
                # All pages blank — return empty with OCR page list.
                # Caller decides whether to enqueue OCR or raise.
                content = ""
            else:
                raise IOError(f"No text extracted from PDF: {path}")

            return content, ocr_needed
        except IOError:
            raise
        except Exception as e:
            raise IOError(f"Failed to extract text from PDF {path}: {e}")

    @staticmethod
    def _is_low_quality_text(text: str, num_pages: int) -> bool:
        """Check if extracted PDF text looks like garbled/schematic output.

        Returns True if the text has very few real words relative to its
        length, suggesting it's labels from a schematic, garbled encoding,
        or minimal embedded text that would benefit from OCR.
        """
        if not text:
            return True
        # Whitespace-separated tokens — normal prose averages 4-6 chars/word.
        # Schematics, garbled PDFs, and label dumps produce long concatenated
        # tokens (e.g. "BCDDate:2025-04-10KiCadE.D.A.8.0.8").
        tokens = text.split()
        if not tokens:
            return True
        avg_token_len = sum(len(t) for t in tokens) / len(tokens)
        if avg_token_len > 15:
            return True
        # Very few tokens per page (normal text has 200+ words/page)
        tokens_per_page = len(tokens) / max(num_pages, 1)
        if tokens_per_page < 30:
            return True
        return False

    def _ocr_pdf_pages(
        self, path: Path, page_indices: list[int], extractor=None,
    ) -> list[tuple[int, str]]:
        """Render PDF pages to images and OCR them.

        Args:
            path: Path to the PDF file.
            page_indices: Zero-based page indices to OCR.
            extractor: ContentExtractor to use. Falls back to
                       self._content_extractor if not provided.
        """
        import logging
        import tempfile
        import time
        _log = logging.getLogger(__name__)

        try:
            import pypdfium2 as pdfium
        except ImportError:
            _log.warning("pypdfium2 not installed; cannot OCR PDF pages")
            return []

        if not extractor:
            _log.warning("No content extractor available for OCR")
            return []

        results: list[tuple[int, str]] = []
        start = time.monotonic()
        pdf = pdfium.PdfDocument(str(path))
        tmp_dir = tempfile.mkdtemp(prefix="keep_ocr_")
        try:
            for i in page_indices:
                page = pdf[i]
                bitmap = page.render(scale=2)
                pil_image = bitmap.to_pil()
                tmp_path = Path(tmp_dir) / f"page_{i}.png"
                pil_image.save(str(tmp_path), format="PNG")
                try:
                    text = extractor.extract(
                        str(tmp_path), "image/png"
                    )
                    if text:
                        cleaned = self._clean_ocr_text(text)
                        confidence = self._estimate_ocr_confidence(cleaned)
                        if confidence >= 0.3 and len(cleaned) > 10:
                            results.append((i, cleaned))
                            _log.debug(
                                "OCR page %d: %d chars, confidence=%.2f",
                                i, len(cleaned), confidence,
                            )
                        else:
                            _log.debug(
                                "OCR page %d: rejected (confidence=%.2f, len=%d)",
                                i, confidence, len(cleaned),
                            )
                except Exception as e:
                    _log.warning("OCR failed for page %d: %s", i, e)
        finally:
            pdf.close()
            import shutil
            shutil.rmtree(tmp_dir, ignore_errors=True)

        elapsed = time.monotonic() - start
        _log.info(
            "OCR complete: %d/%d pages extracted in %.1fs",
            len(results), len(page_indices), elapsed,
        )
        return results

    @staticmethod
    def _clean_ocr_text(text: str) -> str:
        """Clean OCR output: remove junk lines, normalize whitespace."""
        lines = text.split("\n")
        cleaned = []
        for line in lines:
            line = line.strip()
            if len(line) < 2:
                continue
            if len(line) > 20 and " " not in line:
                continue  # likely garbage
            if not any(c.isalnum() for c in line):
                continue
            cleaned.append(line)
        return "\n".join(cleaned)

    @staticmethod
    def _estimate_ocr_confidence(text: str) -> float:
        """Estimate OCR quality: ratio of alphanumeric to total chars."""
        if not text:
            return 0.0
        alnum = sum(1 for c in text if c.isalnum())
        return min(1.0, alnum / len(text))

    def _extract_html_text(self, path: Path) -> str:
        """Extract text from HTML file."""
        try:
            html_content = path.read_text(encoding="utf-8")
            return extract_html_text(html_content)
        except ImportError:
            raise IOError(
                f"HTML text extraction requires 'beautifulsoup4' library. "
                f"Install with: pip install beautifulsoup4\n"
                f"Cannot extract text from HTML: {path}"
            )
        except Exception as e:
            raise IOError(f"Failed to extract text from HTML {path}: {e}")

    @staticmethod
    def _detect_email(content: str) -> bool:
        """Check if text content looks like an RFC 822 email message.

        Requires at least 2 recognized email headers before the first blank
        line, to avoid false positives on random text files.
        """
        # Must have a blank line separating headers from body
        if '\n\n' not in content and '\r\n\r\n' not in content:
            return False
        header_block = content.split('\n\n', 1)[0]
        matches = _EMAIL_HEADER_RE.findall(header_block)
        return len(matches) >= 2

    def _extract_email(self, path: Path) -> tuple[str, dict, list[dict] | None]:
        """Extract body text, header tags, and attachments from an RFC 822 email.

        Returns:
            Tuple of (body_text, tags_dict, attachments).
            Tags include from, to, cc, subject, date, message-id.
            Multi-recipient To/Cc headers produce list-valued tags.
            Attachments is a list of dicts with keys: path, filename,
            content_type, size. Each path is a temp file that the caller
            should clean up after processing. None if no attachments.
        """
        import email
        import email.policy
        import email.utils

        raw = path.read_bytes()
        msg = email.message_from_bytes(raw, policy=email.policy.default)

        # Extract body text
        content = ""
        body = msg.get_body(preferencelist=('plain', 'html'))
        if body is not None:
            try:
                payload = body.get_content()
                if body.get_content_type() == 'text/html':
                    payload = extract_html_text(payload)
                content = payload
            except Exception:
                pass
        if not content:
            # Fallback: try to decode as text
            if msg.get_content_type().startswith('text/'):
                try:
                    content = msg.get_content()
                except Exception:
                    pass
        if not content:
            content = f"[Email: {path.name}]"

        # Strip base64-encoded blocks that may be embedded inline
        # (common in Enron corpus where MIME structure is flattened to text/plain)
        content = _strip_base64_blocks(content)

        # Extract attachments from multipart messages
        attachments: list[dict] = []
        if msg.is_multipart():
            att_dir = None
            for part in msg.iter_attachments():
                ct = part.get_content_type()
                filename = part.get_filename() or f"attachment-{len(attachments) + 1}"
                try:
                    payload_bytes = part.get_payload(decode=True)
                except Exception:
                    payload_bytes = None
                if not payload_bytes:
                    continue
                # Save to temp file inside home dir (FileDocumentProvider
                # blocks paths outside home for safety).
                if att_dir is None:
                    cache_dir = Path.home() / ".cache" / "keep" / "email-att"
                    cache_dir.mkdir(parents=True, exist_ok=True)
                    att_dir = Path(tempfile.mkdtemp(prefix="keep_", dir=cache_dir))
                # Sanitize filename for filesystem safety
                safe_name = Path(filename).name  # strip directory components
                att_path = att_dir / safe_name
                att_path.write_bytes(payload_bytes)
                attachments.append({
                    "path": str(att_path),
                    "filename": filename,
                    "content_type": ct,
                    "size": len(payload_bytes),
                })

        # Extract headers as tags
        tags: dict = {}

        # From (edge-tag)
        from_header = msg.get('From', '')
        if from_header:
            addrs = email.utils.getaddresses([from_header])
            from_values = [addr[1].lower() for addr in addrs if addr[1]]
            if len(from_values) == 1:
                tags['from'] = from_values[0]
            elif from_values:
                tags['from'] = from_values

        # To (edge-tag, multi-value)
        to_headers = msg.get_all('To', [])
        if to_headers:
            addrs = email.utils.getaddresses(to_headers)
            to_values = [addr[1].lower() for addr in addrs if addr[1]]
            if len(to_values) == 1:
                tags['to'] = to_values[0]
            elif to_values:
                tags['to'] = to_values

        # Cc (edge-tag, multi-value)
        cc_headers = msg.get_all('Cc', [])
        if cc_headers:
            addrs = email.utils.getaddresses(cc_headers)
            cc_values = [addr[1].lower() for addr in addrs if addr[1]]
            if len(cc_values) == 1:
                tags['cc'] = cc_values[0]
            elif cc_values:
                tags['cc'] = cc_values

        # Bcc (edge-tag, multi-value)
        bcc_headers = msg.get_all('Bcc', [])
        if bcc_headers:
            addrs = email.utils.getaddresses(bcc_headers)
            bcc_values = [addr[1].lower() for addr in addrs if addr[1]]
            if len(bcc_values) == 1:
                tags['bcc'] = bcc_values[0]
            elif bcc_values:
                tags['bcc'] = bcc_values

        # Subject
        subject = msg.get('Subject', '')
        if subject:
            tags['subject'] = subject

        # Date — parse to ISO 8601
        date_header = msg.get('Date', '')
        if date_header:
            try:
                parsed_date = email.utils.parsedate_to_datetime(date_header)
                tags['date'] = parsed_date.isoformat()
            except (ValueError, TypeError):
                tags['date'] = date_header  # preserve raw if unparseable

        # Message-ID
        msg_id = msg.get('Message-ID', '')
        if msg_id:
            tags['message-id'] = msg_id.strip()

        return content, tags or None, attachments or None


def _is_binary_content_type(content_type: str) -> bool:
    """Check if a content type represents binary data that can't be decoded as text."""
    if content_type.startswith(("image/", "audio/", "video/")):
        return True
    binary_types = {
        "application/octet-stream",
        "application/zip",
        "application/gzip",
        "application/x-tar",
        "application/vnd.apple.pages",
        "application/vnd.apple.keynote",
        "application/vnd.apple.numbers",
        "application/msword",
        "application/vnd.ms-excel",
        "application/vnd.ms-powerpoint",
        "application/wasm",
    }
    return content_type in binary_types


# Content types that FileDocumentProvider can extract text from.
# Maps content-type (or URL suffix) → file extension for temp file.
_EXTRACTABLE_TYPES: dict[str, str] = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
}

_EXTRACTABLE_SUFFIXES: dict[str, str] = {
    "pdf": ".pdf",
    "docx": ".docx",
    "pptx": ".pptx",
}


def _is_extractable_binary(content_type: str, url_suffix: str) -> bool:
    """Check if this binary type can be extracted via FileDocumentProvider."""
    return content_type in _EXTRACTABLE_TYPES or url_suffix in _EXTRACTABLE_SUFFIXES


def _extract_via_file_provider(
    data: bytes, uri: str, content_type: str, url_suffix: str,
) -> tuple[str, str]:
    """Write bytes to a temp file and extract text using FileDocumentProvider.

    Uses a temp file inside the user's home directory so that
    FileDocumentProvider's path-traversal guard is satisfied.
    The temp file is deleted after extraction.

    Returns (extracted_content, content_type).
    """
    ext = _EXTRACTABLE_TYPES.get(content_type) or _EXTRACTABLE_SUFFIXES.get(url_suffix, "")

    # Write inside home dir to satisfy FileDocumentProvider's path guard.
    tmp_dir = Path.home() / ".cache" / "keep"
    tmp_dir.mkdir(parents=True, exist_ok=True)

    with tempfile.NamedTemporaryFile(
        suffix=ext, dir=tmp_dir, delete=True,
    ) as tmp:
        tmp.write(data)
        tmp.flush()

        provider = FileDocumentProvider()
        doc = provider.fetch(tmp.name)
        return doc.content, doc.content_type


class HttpDocumentProvider:
    """Fetches documents from HTTP/HTTPS URLs.

    Requires the `requests` library (optional dependency).
    """
    
    def __init__(self, timeout: int = 30, max_size: int = 10_000_000):
        """Initialize.

        Args:
        timeout: Request timeout in seconds
        max_size: Maximum content size in bytes.
        """
        self.timeout = timeout
        self.max_size = max_size
    
    def supports(self, uri: str) -> bool:
        """Check if this is an HTTP(S) URL."""
        return uri.startswith("http://") or uri.startswith("https://")
    
    @staticmethod
    def _is_private_url(uri: str) -> bool:
        """Check if URL targets a private/internal network address.

        Note: DNS resolution here is inherently TOCTOU — the hostname could
        resolve to a different address by the time requests.get() connects.
        Sufficient for CLI use; a hosted service should enforce this at the
        network layer (firewall/VPC rules) rather than relying on client checks.
        """
        from urllib.parse import urlparse
        import ipaddress
        import socket

        parsed = urlparse(uri)
        hostname = parsed.hostname
        if not hostname:
            return True

        # Block known metadata endpoints and localhost
        if hostname in ("metadata.google.internal",):
            return True

        try:
            addr = ipaddress.ip_address(hostname)
            return (addr.is_private or addr.is_loopback or addr.is_link_local
                    or addr.is_reserved or addr.is_unspecified or addr.is_multicast)
        except ValueError:
            pass  # Not an IP literal — resolve it

        try:
            for _, _, _, _, sockaddr in socket.getaddrinfo(hostname, None):
                addr = ipaddress.ip_address(sockaddr[0])
                if (addr.is_private or addr.is_loopback or addr.is_link_local
                        or addr.is_reserved or addr.is_unspecified or addr.is_multicast):
                    return True
        except socket.gaierror:
            pass  # DNS failure will be caught by requests

        return False

    _MAX_REDIRECTS = 5

    def fetch(self, uri: str) -> Document:
        """Fetch content from HTTP URL with text extraction for HTML."""
        if self._is_private_url(uri):
            raise IOError(f"Blocked request to private/internal address: {uri}")

        from keep.providers.http import http_session

        # Follow redirects manually so each hop is validated against SSRF
        target = uri
        for _ in range(self._MAX_REDIRECTS):
            resp = http_session().get(
                target,
                timeout=self.timeout,
                stream=True,
                allow_redirects=False,
            )
            if resp.is_redirect:
                target = resp.headers.get("Location", "")
                if not target.startswith(("http://", "https://")):
                    raise IOError(f"Redirect to unsupported scheme: {target}")
                if self._is_private_url(target):
                    raise IOError(f"Redirect to private/internal address blocked: {target}")
                resp.close()
                continue
            break
        else:
            raise IOError(f"Too many redirects fetching {uri}")

        try:
            with resp:
                resp.raise_for_status()

                # Check declared size
                content_length = resp.headers.get("content-length")
                if content_length:
                    try:
                        if int(content_length) > self.max_size:
                            raise IOError(f"Content too large: {content_length} bytes")
                    except ValueError:
                        pass  # Malformed header — enforce via iter_content below

                # Read content in chunks with enforced size limit
                chunks: list[bytes] = []
                downloaded = 0
                for chunk in resp.iter_content(chunk_size=65536):
                    downloaded += len(chunk)
                    if downloaded > self.max_size:
                        chunks.append(chunk[:self.max_size - (downloaded - len(chunk))])
                        break
                    chunks.append(chunk)
                raw = b"".join(chunks)

                # Get content type
                content_type = resp.headers.get("content-type", "text/plain")
                if ";" in content_type:
                    content_type = content_type.split(";")[0].strip()

                # Detect content type from URL suffix as fallback
                url_suffix = uri.lower().rsplit("?", 1)[0].rsplit(".", 1)[-1]

                # Extract content based on type
                if content_type == "text/html":
                    encoding = resp.encoding or "utf-8"
                    content = raw.decode(encoding, errors="replace")
                    try:
                        content = extract_html_text(content)
                    except ImportError:
                        pass
                elif _is_extractable_binary(content_type, url_suffix):
                    # Binary type with text extraction support —
                    # write to temp file and use FileDocumentProvider
                    content, content_type = _extract_via_file_provider(
                        raw, uri, content_type, url_suffix,
                    )
                elif _is_binary_content_type(content_type):
                    # Binary type we can't extract text from —
                    # return a placeholder instead of decoded garbage
                    size_kb = len(raw) / 1024
                    content = (
                        f"[Remote binary document: {uri}, "
                        f"type={content_type}, {size_kb:.0f}KB]"
                    )
                else:
                    encoding = resp.encoding or "utf-8"
                    content = raw.decode(encoding, errors="replace")

                return Document(
                    uri=uri,
                    content=content,
                    content_type=content_type,
                    metadata={
                        "status_code": resp.status_code,
                        "headers": dict(resp.headers),
                    },
                )
        except requests.RequestException as e:
            raise IOError(f"Failed to fetch {uri}: {e}")


class CompositeDocumentProvider:
    """Combines multiple document providers, delegating to the appropriate one.
    
    This is the default provider used by Keeper.
    """
    
    def __init__(self, providers: list[DocumentProvider] | None = None):
        """Initialize.

        Args:
        providers: List of providers to try. If None, uses defaults.
        """
        if providers is None:
            self._providers = [
                FileDocumentProvider(),
                HttpDocumentProvider(),
            ]
        else:
            self._providers = list(providers)
    
    def supports(self, uri: str) -> bool:
        """Check if any provider supports this URI."""
        return any(p.supports(uri) for p in self._providers)
    
    def fetch(self, uri: str) -> Document:
        """Fetch using the first provider that supports this URI."""
        for provider in self._providers:
            if provider.supports(uri):
                return provider.fetch(uri)
        
        raise ValueError(f"No provider supports URI: {uri}")
    
    def add_provider(self, provider: DocumentProvider) -> None:
        """Add a provider to the list (checked first)."""
        self._providers.insert(0, provider)


# Register providers
_registry = get_registry()
_registry.register_document("file", FileDocumentProvider)
_registry.register_document("http", HttpDocumentProvider)
_registry.register_document("composite", CompositeDocumentProvider)
